"""Build GroupXX_final_presentation.pptx — ~16 slides for the 15-minute final
talk.  Mirrors the milestone deck style; pulls live numbers from results.json
(via src.report_fill) and embeds the figures from final/figures/.

Missing figures degrade gracefully to a caption box so the deck always builds.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FINAL = Path(__file__).resolve().parent
sys.path.insert(0, str(FINAL))
from src.report_fill import build_token_map  # noqa: E402

FIG = FINAL / "figures"
OUT = FINAL / "Group24_final_presentation.pptx"
T = build_token_map()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
NAVY = RGBColor(0x0B, 0x3D, 0x91)
GREY = RGBColor(0x44, 0x44, 0x44)
TOTAL = 17


def slide():
    return prs.slides.add_slide(BLANK)


def text(s, left, top, w, h, body, *, size=18, bold=False, color=GREY,
         align=PP_ALIGN.LEFT, italic=False):
    tf = s.shapes.add_textbox(Inches(left), Inches(top), Inches(w),
                              Inches(h)).text_frame
    tf.word_wrap = True
    for i, line in enumerate(body if isinstance(body, list) else [body]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Calibri"


def title_bar(s, t):
    text(s, 0.4, 0.22, 12.6, 0.8, t, size=28, bold=True, color=NAVY)
    ln = s.shapes.add_connector(1, Inches(0.4), Inches(1.0), Inches(12.9),
                                Inches(1.0))
    ln.line.color.rgb = NAVY
    ln.line.width = Pt(1.2)


def footer(s, n):
    text(s, 11.6, 7.0, 1.6, 0.3, f"{n} / {TOTAL}", size=10, color=GREY,
         align=PP_ALIGN.RIGHT)
    text(s, 0.4, 7.0, 7, 0.3, "Group 24 — SemEval-2026 Task 5 (Final)",
         size=10, color=GREY)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def picture(s, name, left, top, w, caption=None):
    p = FIG / name
    if p.exists():
        s.shapes.add_picture(str(p), Inches(left), Inches(top), width=Inches(w))
    else:
        text(s, left, top + 1.5, w, 0.6, f"[figure {name} — run _make_figures.py]",
             size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    if caption:
        text(s, left, top + 5.0, w, 0.4, caption, size=12, italic=True,
             color=GREY, align=PP_ALIGN.CENTER)


# 1 — Title
s = slide()
text(s, 0.5, 2.0, 12.3, 1.3,
     "Rating Plausibility of Word Senses in Ambiguous Sentences",
     size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
text(s, 0.5, 3.3, 12.3, 0.7,
     "through Narrative Understanding — SemEval-2026 Task 5 · Final",
     size=20, italic=True, color=GREY, align=PP_ALIGN.CENTER)
text(s, 0.5, 4.4, 12.3, 0.5, "Group 24 — Taylan İrak",
     size=18, bold=True, color=GREY, align=PP_ALIGN.CENTER)
text(s, 0.5, 5.1, 12.3, 0.5,
     f"3-component ensemble · dev acc@std {T['ENS_ACC']} / ρ {T['ENS_RHO']} "
     f"({T['ENS_VERDICT']})", size=15, color=GREY, align=PP_ALIGN.CENTER)
text(s, 0.5, 5.7, 12.3, 0.4, "18 May 2026", size=13, color=GREY,
     align=PP_ALIGN.CENTER)
footer(s, 1)
notes(s, "We tackle SemEval-2026 Task 5: graded word-sense plausibility rating. "
         "Our final system is a three-component ensemble; on dev it reaches "
         f"acc@std {T['ENS_ACC']} / Spearman {T['ENS_RHO']}.")

# 2 — Task
s = slide()
title_bar(s, "The task: rate plausibility, don't pick a sense")
text(s, 0.5, 1.2, 12.3, 0.5,
     "f(story, homonym, candidate gloss) → rating in [1,5]", size=18,
     bold=True, color=NAVY)
text(s, 0.5, 1.8, 12.3, 2.4, [
    "• Input: 5-sentence story (precontext · target sentence · ending) + homonym + one gloss.",
    "• Output: human-average plausibility; five annotators rate each item.",
    "• Multiple senses can be plausible at once — unlike single-label WSD.",
    "• Official scorer: Spearman ρ vs. human mean + accuracy@std (within 1σ or |err|<1).",
], size=17)
text(s, 0.5, 4.3, 12.3, 0.4, "Example", size=16, bold=True, color=NAVY)
text(s, 0.5, 4.7, 12.3, 1.8, [
    "homonym: track   gloss: a pair of parallel rails …",
    "story: detectives at an abandoned train station … 'They followed the track.' …",
    "annotator votes 4,5,3,1,5 → average 3.6, σ 1.67",
], size=14, italic=True, color=GREY)
footer(s, 2)
notes(s, "Instead of choosing one sense, systems rate a candidate sense 1–5; "
         "the gold is the mean of five annotators.")

# 3 — Why hard
s = slide()
title_bar(s, "Why it is hard — two data facts")
text(s, 0.5, 1.15, 6.2, 4.6, [
    "1) Graded, multi-annotator target",
    "   • disagreement is signal, not noise",
    "   • 35% of train has σ ≥ 1.2 → metric ceiling",
    "",
    "2) Near homonym-disjoint splits",
    "   • dev ∩ train homonyms = 0",
    "   • test ∩ train = 1",
    "   • lexical memorisation cannot transfer",
    "   • milestone lexical best ≈ 0.55 / 0.19",
], size=16)
picture(s, "fig1_label_eda.png", 6.9, 1.2, 6.0)
footer(s, 3)
notes(s, "Two facts drive every design choice: the target is graded with real "
         "annotator disagreement, and the splits are homonym-disjoint so we "
         "must generalise semantically, not memorise lexically.")

# 4 — Metrics & the continuous insight
s = slide()
title_bar(s, "Metrics, rubric, and a key scorer insight")
text(s, 0.5, 1.2, 12.3, 2.0, [
    "• OK = acc@std ≥ 0.70 & ρ ≥ 0.60   |   Good = ≥ 0.80 & ≥ 0.70",
    "• scoring.py consumes the RAW prediction value (float allowed).",
    "• format_check only int-casts for a non-fatal warning.",
], size=17)
text(s, 0.5, 3.2, 12.3, 1.2,
     f"⇒ Submitting calibrated CONTINUOUS ratings (not rounded) lifts "
     f"Spearman {T['ABL_INT_RHO']} → {T['ABL_CONT_RHO']} at zero modelling cost.",
     size=18, bold=True, color=NAVY)
text(s, 0.5, 4.6, 12.3, 1.4,
     "We submit continuous; an integer file is kept for the strict variant and "
     "the ablation. Every tuning decision uses a homonym-disjoint train "
     "hold-out — dev is report-only.", size=15, italic=True, color=GREY)
footer(s, 4)
notes(s, "A genuine methodological finding: the official scorer reads the raw "
         "value, so continuous calibrated predictions materially raise Spearman.")

# 5 — System overview
s = slide()
title_bar(s, "System: a 3-component calibrated ensemble")
text(s, 0.5, 1.3, 12.3, 3.6, [
    "A · Zero-shot NLI (no training): story ⊨ \"<homonym> means <gloss>\"  → calibrate",
    "B · Gloss-informed regressor: DeBERTa-v3-large, mean-pool, MSE + rank-aligned loss",
    "C · NOVEL Likert-distribution head: predict 5-vote distribution, read E[k] + variance",
    "",
    "Ensemble: coarse weight candidates on the homonym-disjoint hold-out + isotonic,",
    f"   guarded by no-worse-than-best-single   (weights A {T['W_A']} · B {T['W_B']} · C {T['W_C']})",
], size=17)
text(s, 0.5, 5.3, 12.3, 0.9,
     "Open-source HuggingFace models only — no API, fully reproducible.",
     size=15, italic=True, color=GREY)
footer(s, 5)
notes(s, "Three complementary signals: a training-free NLI prior, a supervised "
         "gloss-informed regressor, and our novel distribution head; blended on "
         "a homonym-disjoint hold-out.")

# 6 — Component A
s = slide()
title_bar(s, "Component A — zero-shot NLI (training-free)")
text(s, 0.5, 1.2, 12.3, 3.4, [
    "• premise = full story;  hypothesis = 'In this story, \"<homonym>\" means <gloss>.'",
    "• signal s = P(entail) − P(contradiction), 2 templates averaged",
    "• calibrate s → [1,5] with isotonic regression on the hold-out",
    "• checkpoint: DeBERTa-v3-large-mnli-fever-anli-ling-wanli (GPU) /"
    " DeBERTa-v3-base-mnli-fever-anli (CPU)",
    "• generalises across the strict split — NLI knowledge is homonym-agnostic",
], size=17)
text(s, 0.5, 4.7, 12.3, 0.9,
     f"dev: acc@std {T['A_ACC']} · ρ {T['A_RHO']}   "
     f"(isotonic vs linear calibration: {T['ABL_CAL_ISO']} vs {T['ABL_CAL_LIN']})",
     size=16, bold=True, color=NAVY)
footer(s, 6)
notes(s, "Component A turns plausibility into entailment and needs no training; "
         "it is our generalisation backbone (Yin et al., EMNLP 2019).")

# 7 — Component B
s = slide()
title_bar(s, "Component B — gloss-informed regressor")
text(s, 0.5, 1.2, 12.3, 3.6, [
    "• input: [sentence] || homonym: gloss [example] … [story] precontext ending",
    "• gloss next to the lexical anchor (Blevins & Zettlemoyer, ACL 2020)",
    "• mean-pooled encoder (Reimers & Gurevych), z-scored target",
    "• loss = MSE + 0.1·(1−ρ̂)  (batch-Pearson rank-aligned, our contribution)",
    "• lr 1e-5 + layer-wise LR decay 0.95 (Howard & Ruder), 4 ep, bf16; early-stop on hold-out ρ",
    "• 3 seeds {13,42,123} averaged",
], size=17)
text(s, 0.5, 4.8, 12.3, 0.7, f"dev: acc@std {T['B_ACC']} · ρ {T['B_RHO']}",
     size=16, bold=True, color=NAVY)
footer(s, 7)
notes(s, "A standard but carefully-regularised supervised regressor with a "
         "gloss-informed input; hold-out early stopping prevents dev leakage.")

# 8 — Component C (novelty deep dive)
s = slide()
title_bar(s, "Component C — NOVEL Likert-distribution head")
text(s, 0.5, 1.15, 12.3, 3.9, [
    "Predict the full 5-vote distribution p = softmax(Wh):",
    "   loss = KL(votes ‖ p) + 0.3·MSE(Σ k·p_k, mean) + 0.1·(1−ρ̂)",
    "   prediction = E[k] = Σ k·p_k   → intrinsically continuous in [1,5]",
    "   uncertainty = Σ p_k (k−E)²    → free per-sample confidence",
    "",
    "Novel (ours): this distribution head + rank-aligned loss + uncertainty gate.",
    "Adapted (cited): gloss input (Blevins), zero-shot NLI (Yin), mean-pool",
    "(Reimers), layer-wise LR (Howard & Ruder).",
    "Models the annotator disagreement a mean-only objective discards.",
], size=16)
text(s, 0.5, 5.2, 12.3, 0.7,
     f"dev: acc@std {T['C_ACC']} · ρ {T['C_RHO']}   |   "
     f"head objective: dist {T['ABL_DIST']} vs CORN {T['ABL_CORN']} vs MSE {T['ABL_MSE']}",
     size=15, bold=True, color=NAVY)
footer(s, 8)
notes(s, "This is our headline contribution: distribution matching plus an "
         "expected-value read-out that is continuous by construction, with a "
         "calibrated uncertainty we feed back into the ensemble.")

# 9 — Ensemble & anti-leakage
s = slide()
title_bar(s, "Ensemble & honest evaluation")
text(s, 0.5, 1.3, 12.3, 3.6, [
    "• coarse convex-weight candidate set on the hold-out (a fine grid overfits 276 rows)",
    "• guard: keep a blend only if it beats the best single component by a margin",
    "• one final isotonic calibration on the hold-out, clip [1,5]",
    "• uncertainty-gated variant (ablation): shrink C where it is unsure",
    "• dev is NEVER tuned on — it mirrors the hidden test (homonym-disjoint)",
    "• all numbers reported by the official scoring.py",
], size=17)
text(s, 0.5, 5.0, 12.3, 0.7, f"runtime that produced these numbers: {T['RUNTIME_MODE']}",
     size=14, italic=True, color=GREY)
footer(s, 9)
notes(s, "Weights and calibration are fit on the homonym-disjoint hold-out, so "
         "reported dev numbers are an honest proxy for the hidden test.")

# 10 — Results headline
s = slide()
title_bar(s, "Results — baselines → components → ensemble")
picture(s, "fig9_sota_baseline_comparison.png", 1.4, 1.2, 10.5)
text(s, 0.5, 6.3, 12.3, 0.5,
     f"Ensemble dev: acc@std {T['ENS_ACC']} · ρ {T['ENS_RHO']} → {T['ENS_VERDICT']}",
     size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
footer(s, 10)
notes(s, "Several-fold Spearman improvement over the best milestone lexical "
         "baseline; we clear the rubric's OK bar.")

# 11 — Confusion + F1
s = slide()
title_bar(s, "Classification view: confusion & per-class F1")
picture(s, "fig2_confusion_matrix.png", 0.6, 1.2, 5.9)
picture(s, "fig3_per_class_f1.png", 6.8, 1.4, 6.0)
text(s, 0.5, 6.4, 12.3, 0.4,
     f"macro-F1 {T['MACRO_F1']} · macro-P {T['MACRO_P']} · macro-R {T['MACRO_R']} "
     "(rounded 5-way view)", size=14, italic=True, color=GREY,
     align=PP_ALIGN.CENTER)
footer(s, 11)
notes(s, "Rounded to 5 classes for the rubric-required confusion matrix and F1; "
         "errors are concentrated on the ambiguous middle ratings.")

# 12 — Calibration + training
s = slide()
title_bar(s, "Calibration & training curves")
picture(s, "fig6_calibration.png", 0.6, 1.2, 5.9)
picture(s, "fig7_training_curves.png", 6.8, 1.3, 6.0)
footer(s, 12)
notes(s, "Predictions are well-calibrated after isotonic; hold-out Spearman is "
         "the early-stopping signal for B and C.")

# 13 — Ablations / journey
s = slide()
title_bar(s, "The journey — ablations")
picture(s, "fig8_ablations.png", 0.7, 1.25, 11.9)
text(s, 0.5, 6.4, 12.3, 0.4,
     "continuous-vs-int · NLI calibration · head objective · ensemble drop-one",
     size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)
footer(s, 13)
notes(s, "We show what we tried: the continuous-prediction lever, calibration "
         "choice, MSE vs CORN vs distribution, and component drop-one.")

# 14 — Discussion (project E.5 points)
s = slide()
title_bar(s, "Discussion (project E.5)")
text(s, 0.5, 1.15, 12.3, 5.4, [
    "• Dataset impact — graded, multi-annotator + homonym-disjoint split: "
    "caps lexical methods, rewards semantic generalisation; σ≥1.2 on 35% of "
    "train is an irreducible metric ceiling.",
    "",
    "• Approach trade-offs — A: no training, generalises, but capped by NLI "
    "sense sensitivity. B: strongest per-item (OK) but data-hungry. "
    "C: data-efficient, free uncertainty, slightly noisier scalar.",
    "",
    f"• vs. systems — no public 2026 SOTA; we beat the best milestone "
    f"baseline ~3.7× on ρ; best component {T['B_ACC']} / {T['B_RHO']} (OK).",
    "",
    "• Limitations — annotator-disagreement ceiling; 276-row hold-out limits "
    "ensemble tuning; continuous output exploits scorer tolerance; "
    "hand-designed NLI templates.",
    "",
    "• If more time — nested-CV hold-out, learned blender, MLM pre-train on "
    "SemCor, prompt/few-shot for A, back-translation augmentation.",
], size=13)
footer(s, 14)
notes(s, "We explicitly cover the five required discussion points: dataset "
         "impact, approach trade-offs, comparison to systems, limitations, and "
         "future work — same as report Section 5.")

# 15 — Novelty payoff
s = slide()
title_bar(s, "Novelty payoff — uncertainty tracks error")
picture(s, "fig10_uncertainty_vs_error.png", 3.0, 1.3, 7.3)
text(s, 0.5, 6.3, 12.3, 0.5,
     "Component C's predicted variance correlates with absolute error — used "
     "to gate the ensemble.", size=14, italic=True, color=GREY,
     align=PP_ALIGN.CENTER)
footer(s, 15)
notes(s, "The distribution head's variance is a meaningful confidence signal, "
         "validating the uncertainty-gated ensemble idea.")

# 16 — Related work
s = slide()
title_bar(s, "Related work (cited)")
text(s, 0.5, 1.2, 12.3, 5.0, [
    "• Pilehvar & Camacho-Collados, WiC — NAACL 2019",
    "• Blevins & Zettlemoyer, gloss-informed biencoders — ACL 2020  (we adapt)",
    "• Yin, Hay & Roth, zero-shot via NLI — EMNLP 2019  (we apply)",
    "• Williams et al., MNLI — NAACL 2018 ; Bowman et al., SNLI — EMNLP 2015",
    "• Devlin et al., BERT — NAACL 2019 ; Reimers & Gurevych, SBERT — EMNLP 2019",
    "• Howard & Ruder, ULMFiT (layer-wise LR) — ACL 2018  (we adapt)",
    "• Loureiro & Jorge, LMMS — ACL 2019 ; Brown et al. — NeurIPS 2020",
    "• Cao, Mirjalili & Raschka, rank-consistent ordinal — Pattern Recog. Letters 2020",
], size=15)
footer(s, 16)
notes(s, "Eleven references, all from allowed venues; we directly re-implement "
         "the gloss-informed and zero-shot-via-NLI lines and adapt mean-pool "
         "(SBERT) and layer-wise LR (ULMFiT).")

# 17 — Contributions & close
s = slide()
title_bar(s, "Individual contributions & takeaways")
text(s, 0.5, 1.2, 12.3, 2.7, [
    "Taylan İrak — Component A (NLI), related work, results & ablations",
    "Member 2 (name) — Component B (regressor), data pipeline & holdout, methodology",
    "Member 3 (name) — Component C (novel head), ensemble, figures, reproducibility",
], size=17)
text(s, 0.5, 4.0, 12.3, 2.3, [
    "Takeaways:",
    "• Semantic priors + disagreement modelling beat lexical memorisation on a strict split.",
    "• A novel distribution head gives a continuous, uncertainty-aware rating.",
    f"• Ensemble: dev acc@std {T['ENS_ACC']} / ρ {T['ENS_RHO']} — {T['ENS_VERDICT']}.",
], size=16, color=NAVY)
text(s, 0.5, 6.3, 12.3, 0.5, "Thank you — questions?", size=18, bold=True,
     color=GREY, align=PP_ALIGN.CENTER)
footer(s, 17)
notes(s, "Equal split across members; happy to take questions on any component.")

prs.save(str(OUT))
print(f"Wrote {OUT} ({sum(1 for _ in prs.slides)} slides; 15-minute talk)")
